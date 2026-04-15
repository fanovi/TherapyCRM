#!/usr/bin/env python3
"""
Script per generare la presentazione PowerPoint di TherapyCRM
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colori del tema
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)      # Blu scuro
SECONDARY = RGBColor(0x3B, 0x82, 0xF6)    # Blu medio
ACCENT = RGBColor(0x06, 0xB6, 0xD4)       # Ciano
DARK = RGBColor(0x1F, 0x29, 0x37)         # Grigio scuro
GRAY = RGBColor(0x64, 0x74, 0x8B)         # Grigio
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)     # Sfondo chiaro
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)        # Verde
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)       # Arancione
RED = RGBColor(0xEF, 0x44, 0x44)          # Rosso
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)       # Viola

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color):
    """Aggiunge sfondo colorato alla slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Aggiunge un rettangolo colorato"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Aggiunge un textbox"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, bullet=False, font_name='Calibri'):
    """Aggiunge un paragrafo"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet:
        p.level = 0
    return p


def add_card(slide, left, top, width, height, title, items, icon_color=PRIMARY):
    """Aggiunge una card con titolo e lista"""
    # Sfondo card
    card = add_shape_bg(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False

    # Barra colorata in alto
    add_shape_bg(slide, left, top, width, Inches(0.06), icon_color)

    # Titolo card
    add_textbox(slide, left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.5),
                title, font_size=16, color=icon_color, bold=True)

    # Items
    txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.7), width - Inches(0.5), height - Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.font.name = 'Calibri'
        p.space_before = Pt(3)
        p.space_after = Pt(3)
    return card


def add_footer(slide, slide_num, total):
    """Aggiunge footer con numero slide"""
    add_shape_bg(slide, Inches(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), PRIMARY)
    add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.35), Inches(3), Inches(0.3),
                "TherapyCRM - Piattaforma Gestionale per Centri Terapeutici",
                font_size=9, color=WHITE)
    add_textbox(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.35), Inches(1), Inches(0.3),
                f"{slide_num} / {total}", font_size=9, color=WHITE, alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 16

# =============================================================================
# SLIDE 1 - COPERTINA
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, PRIMARY)

# Decorazione grafica
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), ACCENT)
add_shape_bg(slide, Inches(0), SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), ACCENT)

# Cerchio decorativo
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, SLIDE_W - Inches(4), Inches(-1), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x25, 0x4B, 0xBF)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), SLIDE_H - Inches(3), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x25, 0x4B, 0xBF)
circle2.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.8), Inches(9), Inches(1.2),
            "TherapyCRM", font_size=54, color=WHITE, bold=True)

add_textbox(slide, Inches(1), Inches(3.0), Inches(9), Inches(0.8),
            "Piattaforma Integrata per la Gestione di Centri Terapeutici",
            font_size=26, color=RGBColor(0xBF, 0xDB, 0xFE))

# Linea separatore
add_shape_bg(slide, Inches(1), Inches(4.0), Inches(3), Inches(0.05), ACCENT)

add_textbox(slide, Inches(1), Inches(4.3), Inches(9), Inches(0.6),
            "Gestionale Web  |  App Mobile  |  Pannello Amministrativo",
            font_size=18, color=RGBColor(0x93, 0xC5, 0xFD))

add_textbox(slide, Inches(1), Inches(5.8), Inches(5), Inches(0.4),
            "Soluzione completa per centri riabilitativi e terapeutici multi-sede",
            font_size=14, color=RGBColor(0x93, 0xC5, 0xFD))


# =============================================================================
# SLIDE 2 - IL PROBLEMA
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
            "Le Sfide della Gestione di un Centro Terapeutico",
            font_size=30, color=WHITE, bold=True)

problems = [
    ("Gestione Appuntamenti Complessa",
     "Coordinare decine di terapisti, centinaia di pazienti e migliaia di appuntamenti settimanali con strumenti manuali o fogli Excel."),
    ("Mancanza di Visibilita in Tempo Reale",
     "Assenze, sovrapposizioni e conflitti di calendario scoperti troppo tardi, causando disservizi e inefficienze."),
    ("Comunicazione Frammentata",
     "Informazioni disperse tra email, telefonate e messaggi, senza un canale unico e tracciabile."),
    ("Assenza di Dati e Statistiche",
     "Impossibilita di analizzare trend, monitorare performance e prendere decisioni basate su dati concreti."),
    ("Gestione Documentale Inefficiente",
     "Richieste di certificati e documenti gestite manualmente, con tempi di evasione lunghi e nessun tracciamento."),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.6) + Inches(i * 1.1)
    # Numero
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.5), Inches(0.5))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = RED if i < 2 else ORANGE if i < 4 else SECONDARY
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    add_textbox(slide, Inches(1.5), y - Inches(0.05), Inches(10), Inches(0.35),
                title, font_size=18, color=DARK, bold=True)
    add_textbox(slide, Inches(1.5), y + Inches(0.3), Inches(10), Inches(0.5),
                desc, font_size=13, color=GRAY)

add_footer(slide, 2, TOTAL_SLIDES)


# =============================================================================
# SLIDE 3 - LA SOLUZIONE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
            "La Soluzione: TherapyCRM", font_size=30, color=WHITE, bold=True)

# Tre pilastri
pillars = [
    ("Gestionale Web", "Pannello di controllo completo per manager e coordinatori. Gestione centralizzata di pazienti, terapisti, piani terapeutici e calendario.", SECONDARY),
    ("App Mobile", "Applicazione dedicata per terapisti e famiglie dei pazienti. Calendario, segnalazione assenze, richieste documenti e reclami.", GREEN),
    ("Pannello Super Admin", "Controllo totale del sistema: gestione utenti, ruoli, permessi, distretti sanitari e log attivita.", PURPLE),
]

for i, (title, desc, color) in enumerate(pillars):
    x = Inches(0.5) + Inches(i * 4.2)
    # Card
    card = add_shape_bg(slide, x, Inches(1.6), Inches(3.9), Inches(3.5), LIGHT_BG)
    # Barra superiore
    add_shape_bg(slide, x, Inches(1.6), Inches(3.9), Inches(0.08), color)
    # Icona cerchio
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.45), Inches(2.0), Inches(1), Inches(1))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()
    tf = icon.text_frame
    icons = ["WEB", "APP", "ADM"]
    tf.paragraphs[0].text = icons[i]
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(0.5),
                title, font_size=20, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.7), Inches(3.5), Inches(1.2),
                desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Freccia/connessione
add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.8),
            "Un ecosistema integrato dove ogni componente comunica in tempo reale con gli altri, garantendo coerenza dei dati e massima efficienza operativa.",
            font_size=14, color=DARK, alignment=PP_ALIGN.CENTER)

add_footer(slide, 3, TOTAL_SLIDES)


# =============================================================================
# SLIDE 4 - GESTIONALE: DASHBOARD
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Dashboard e Panoramica Generale", font_size=26, color=WHITE, bold=True)

features = [
    "Panoramica immediata di tutti i KPI del centro",
    "Conteggio pazienti attivi, terapisti operativi, appuntamenti in programma",
    "Notifiche e comunicazioni recenti in evidenza",
    "Accesso rapido a tutte le sezioni del gestionale",
    "Interfaccia moderna e responsive (desktop, tablet)",
    "Supporto modalita chiara e scura (Dark Mode)",
    "Sidebar navigabile e collassabile per massima usabilita",
]

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
for i, feat in enumerate(features):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"  {feat}"
    p.font.size = Pt(15)
    p.font.color.rgb = DARK
    p.font.name = 'Calibri'
    p.space_before = Pt(8)
    p.space_after = Pt(4)
    # Bullet point
    run = p.runs[0] if p.runs else p.add_run()

# Mock dashboard box
mock = add_shape_bg(slide, Inches(7.5), Inches(1.5), Inches(5.3), Inches(4.5), LIGHT_BG)

# Mini cards nella dashboard mock
mini_cards = [
    ("Pazienti Attivi", "248", GREEN),
    ("Terapisti", "32", SECONDARY),
    ("Appuntamenti Oggi", "87", ORANGE),
    ("Piani Attivi", "156", PURPLE),
]
for i, (label, val, col) in enumerate(mini_cards):
    x = Inches(7.8) + Inches((i % 2) * 2.5)
    y = Inches(1.8) + Inches((i // 2) * 1.5)
    c = add_shape_bg(slide, x, y, Inches(2.2), Inches(1.2), WHITE)
    add_shape_bg(slide, x, y, Inches(0.08), Inches(1.2), col)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(1.8), Inches(0.3),
                label, font_size=10, color=GRAY)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.5), Inches(1.8), Inches(0.5),
                val, font_size=28, color=col, bold=True)

add_textbox(slide, Inches(7.8), Inches(5.0), Inches(5), Inches(0.5),
            "* Esempio illustrativo della dashboard", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

add_footer(slide, 4, TOTAL_SLIDES)


# =============================================================================
# SLIDE 5 - GESTIONALE: GESTIONE PAZIENTI
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Gestione Completa dei Pazienti", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2),
         "Anagrafica Paziente", [
             "  Dati anagrafici completi",
             "  Codice fiscale (generazione automatica)",
             "  Comune e provincia di residenza",
             "  Distretto sanitario (ASL)",
             "  Regime sanitario",
             "  Stato attivo/inattivo",
         ], SECONDARY)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.2),
         "Account Familiari", [
             "  Creazione account per familiari/tutori",
             "  Tipi: genitore, tutore, caregiver",
             "  Collegamento paziente-account",
             "  Credenziali di accesso all'App",
             "  Download PDF credenziali",
             "  Reset password da gestionale",
         ], GREEN)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.2),
         "Funzionalita Avanzate", [
             "  Ricerca avanzata con filtri multipli",
             "  Storico piani terapeutici",
             "  Visualizzazione calendario dedicato",
             "  Gestione documenti e richieste",
             "  Notifiche push verso familiari",
             "  Export dati",
         ], PURPLE)

add_footer(slide, 5, TOTAL_SLIDES)


# =============================================================================
# SLIDE 6 - GESTIONALE: TERAPISTI
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Gestione Terapisti e Team", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2),
         "Profilo Terapista", [
             "  Dati anagrafici e professionali",
             "  Specializzazione (Logopedia, ABA, ecc.)",
             "  Tipo: Dipendente o Consulente",
             "  Ore settimanali contrattuali",
             "  Stato attivo/inattivo",
             "  Credenziali accesso App",
         ], SECONDARY)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.2),
         "Gruppi Coordinatori", [
             "  Creazione gruppi di lavoro",
             "  Assegnazione coordinatore responsabile",
             "  Associazione terapisti al gruppo",
             "  Visibilita filtrata per coordinatore",
             "  Gestione presenze per reparto",
             "  Sostituzione terapisti",
         ], GREEN)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.2),
         "Controllo Operativo", [
             "  Monitoraggio ore settimanali lavorate",
             "  Alert superamento ore contrattuali",
             "  Calendario dedicato per terapista",
             "  Storico assenze con motivazione",
             "  Export report presenze in Excel",
             "  Gestione sostituzioni",
         ], ORANGE)

add_footer(slide, 6, TOTAL_SLIDES)


# =============================================================================
# SLIDE 7 - GESTIONALE: PIANI TERAPEUTICI
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Piani Terapeutici Intelligenti", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.4),
         "Creazione e Configurazione", [
             "  Piano terapeutico associato al paziente con data inizio e fine",
             "  Assegnazione distretto sanitario (ASL) e regime",
             "  Aggiunta di piu terapie per piano (es. Logopedia + Psicomotricita)",
             "  Assegnazione terapista specifico per ogni terapia",
         ], SECONDARY)

add_card(slide, Inches(6.6), Inches(1.5), Inches(6.2), Inches(2.4),
         "Gestione Appuntamenti Automatica", [
             "  Creazione pattern ricorrenti (es. ogni lunedi alle 10:00)",
             "  Generazione automatica appuntamenti per il periodo del piano",
             "  Rilevamento automatico conflitti terapista e paziente",
             "  Controllo duplicati stessa terapia nello stesso giorno",
         ], GREEN)

add_card(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.4),
         "Validazioni Intelligenti", [
             "  Conflitto terapista: nessuna sovrapposizione oraria",
             "  Conflitto paziente: nessuna sovrapposizione per lo stesso paziente",
             "  Limite settimanale: alert se il terapista supera le ore contrattuali",
             "  Conflitto tipologia: max 1 trattamento dello stesso tipo al giorno",
         ], ORANGE)

add_card(slide, Inches(6.6), Inches(4.2), Inches(6.2), Inches(2.4),
         "Monitoraggio e Scadenze", [
             "  Dashboard piani in scadenza con alert automatici",
             "  Tassi di completamento per piano e paziente",
             "  Storico completo delle modifiche (audit trail)",
             "  Notifiche automatiche ai familiari per nuovi appuntamenti",
         ], PURPLE)

add_footer(slide, 7, TOTAL_SLIDES)


# =============================================================================
# SLIDE 8 - GESTIONALE: CALENDARIO
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Calendario Interattivo", font_size=26, color=WHITE, bold=True)

features_left = [
    ("Vista Completa", [
        "  Visualizzazione giornaliera, settimanale e mensile",
        "  Tutti gli appuntamenti di tutti i terapisti in un colpo d'occhio",
        "  Codifica colore per tipologia di trattamento",
        "  Filtro per singolo terapista o paziente",
    ]),
    ("Gestione Appuntamenti", [
        "  Creazione rapida appuntamenti dal calendario",
        "  Modifica drag & drop degli orari",
        "  Cancellazione con conferma e notifica",
        "  Aggiornamento massivo pattern ricorrenti",
    ]),
]

features_right = [
    ("Assenze e Disponibilita", [
        "  Visualizzazione assenze terapisti sul calendario",
        "  Segnalazione assenze pazienti da App",
        "  Gestione assenze con motivazione",
        "  Ricalcolo automatico disponibilita",
    ]),
    ("Impostazioni Calendario", [
        "  Range temporale configurabile (passato/futuro)",
        "  Configurazione orari apertura centro",
        "  Personalizzazione per sede/distretto",
        "  Integrazione React per massime performance",
    ]),
]

for col, features in enumerate([features_left, features_right]):
    x_base = Inches(0.5) + Inches(col * 6.3)
    y_pos = Inches(1.5)
    for title, items in features:
        add_card(slide, x_base, y_pos, Inches(6), Inches(2.4), title, items,
                 SECONDARY if col == 0 else GREEN)
        y_pos += Inches(2.7)

add_footer(slide, 8, TOTAL_SLIDES)


# =============================================================================
# SLIDE 9 - GESTIONALE: STATISTICHE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Modulo Statistiche e Analisi Dati", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.4),
         "Dashboard Statistiche", [
             "  Cards riassuntive con KPI chiave",
             "  Grafici dinamici interattivi",
             "  Auto-refresh ogni 5 minuti",
             "  Links rapidi alle analisi dettagliate",
         ], SECONDARY)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(2.4),
         "Analisi Assenze", [
             "  Heatmap interattiva (orari x giorni)",
             "  Statistiche per motivo assenza",
             "  Trend temporali e tassi mensili",
             "  Filtri per terapista e trattamento",
         ], RED)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(2.4),
         "Analisi Pazienti", [
             "  Demografia: eta, genere, distribuzione",
             "  Trattamenti multipli e combinazioni",
             "  Filtri per regime sanitario e distretto",
             "  Crescita nel tempo e comparazione",
         ], GREEN)

add_card(slide, Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.4),
         "Analisi Trattamenti", [
             "  Ranking per popolarita e pazienti",
             "  Combinazioni trattamenti piu frequenti",
             "  Distribuzione ore settimanali",
             "  Analisi efficacia e completamento",
         ], ORANGE)

add_card(slide, Inches(4.6), Inches(4.2), Inches(3.8), Inches(2.4),
         "Analisi Piani Terapeutici", [
             "  Stati e durate con statistiche aggregate",
             "  Alert piani in scadenza (<=7 giorni)",
             "  Tassi completamento per paziente",
             "  Trend creazione mensile",
         ], PURPLE)

add_card(slide, Inches(8.7), Inches(4.2), Inches(4.1), Inches(2.4),
         "Export e Reportistica", [
             "  Export completo in formato Excel",
             "  Report filtrabili per periodo e parametri",
             "  Dati ottimizzati con caching intelligente",
             "  Grafici stampabili per riunioni e report",
         ], DARK)

add_footer(slide, 9, TOTAL_SLIDES)


# =============================================================================
# SLIDE 10 - GESTIONALE: COMUNICAZIONI E NOTIFICHE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Comunicazioni, Documenti e Reclami", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2),
         "Comunicazioni Interne", [
             "  Sistema messaggistica integrato",
             "  Notifiche tra operatori del gestionale",
             "  Filtro per tipo: tutte, interne, non lette",
             "  Segna come letto singola/tutte",
             "  Storico completo comunicazioni",
             "  Notifiche push verso App mobile",
         ], SECONDARY)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.2),
         "Gestione Documenti e Richieste", [
             "  Richieste documenti da familiari via App",
             "  Tipologie: certificati, relazioni, attestati",
             "  Workflow stati: in attesa, accettata,",
             "    in lavorazione, pronta, consegnata",
             "  Data stimata di completamento",
             "  Download documenti completati",
             "  Tracciamento completo del ciclo di vita",
         ], GREEN)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.2),
         "Gestione Reclami", [
             "  Ricezione reclami da App familiari",
             "  Titolo, descrizione e paziente associato",
             "  Limite: max 1 reclamo per paziente/giorno",
             "  Visualizzazione e gestione dal gestionale",
             "  Storico reclami per paziente",
             "  Tracciabilita completa",
         ], RED)

add_footer(slide, 10, TOTAL_SLIDES)


# =============================================================================
# SLIDE 11 - GESTIONALE: UTENTI E RUOLI
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SECONDARY)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "GESTIONALE WEB", font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Gestione Utenti e Permessi Granulari", font_size=26, color=WHITE, bold=True)

# Ruoli
roles = [
    ("Super Admin", "Accesso completo al sistema, gestione infrastruttura, ruoli e configurazioni globali", DARK),
    ("Amministratore", "Gestione completa di pazienti, terapisti, piani terapeutici e impostazioni del centro", PRIMARY),
    ("Manager", "Supervisione operativa, accesso a statistiche, gestione comunicazioni e coordinamento", SECONDARY),
    ("Coordinatore", "Visibilita sul proprio gruppo di terapisti, gestione presenze e assenze del reparto", GREEN),
    ("Terapista (App)", "Accesso al proprio calendario, segnalazione assenze, visualizzazione appuntamenti", ORANGE),
    ("Familiare/Paziente (App)", "Visualizzazione calendario, richiesta documenti, segnalazione assenze, reclami", PURPLE),
]

for i, (role, desc, color) in enumerate(roles):
    y = Inches(1.5) + Inches(i * 0.9)
    # Badge ruolo
    badge = add_shape_bg(slide, Inches(0.8), y, Inches(2.5), Inches(0.45), color)
    tf = badge.text_frame
    tf.paragraphs[0].text = role
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = 'Calibri'

    add_textbox(slide, Inches(3.5), y, Inches(9), Inches(0.5),
                desc, font_size=13, color=DARK)

# Box permessi
perm_box = add_shape_bg(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), LIGHT_BG)
add_textbox(slide, Inches(1.0), Inches(6.05), Inches(11.5), Inches(0.6),
            "Oltre 30 permessi granulari configurabili: view_calendar, manage_appointments, view_statistics, export_data, manage_documents, view_notifications e molti altri",
            font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

add_footer(slide, 11, TOTAL_SLIDES)


# =============================================================================
# SLIDE 12 - APP MOBILE: TERAPISTI
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "APP MOBILE", font_size=14, color=RGBColor(0xA7, 0xF3, 0xD0))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "L'App per i Terapisti", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2),
         "Calendario Personale", [
             "  Visualizzazione appuntamenti giornalieri",
             "  Dettagli paziente e tipo trattamento",
             "  Orario e durata di ogni sessione",
             "  Navigazione per data",
             "  Sincronizzazione in tempo reale",
         ], GREEN)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.2),
         "Gestione Assenze", [
             "  Segnalazione assenza paziente",
             "  Finestra temporale: da +1 a +15 giorni",
             "  Motivazione assenza obbligatoria",
             "  Rimozione assenza (fino a 1 ora prima)",
             "  Notifica automatica al gestionale",
         ], ORANGE)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.2),
         "Sicurezza e Accesso", [
             "  Autenticazione JWT sicura",
             "  Login biometrico (impronta/Face ID)",
             "  Cambio password obbligatorio ogni 90 giorni",
             "  Refresh token automatico (30gg)",
             "  Notifiche push in tempo reale",
         ], SECONDARY)

add_footer(slide, 12, TOTAL_SLIDES)


# =============================================================================
# SLIDE 13 - APP MOBILE: FAMILIARI/PAZIENTI
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "APP MOBILE", font_size=14, color=RGBColor(0xA7, 0xF3, 0xD0))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "L'App per Familiari e Pazienti", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2),
         "Calendario Appuntamenti", [
             "  Visualizzazione tutti gli appuntamenti",
             "  Dettagli: terapista, tipo trattamento, orario",
             "  Navigazione intuitiva per data",
             "  Stato appuntamento in tempo reale",
             "  Accesso per ogni paziente collegato",
         ], GREEN)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.2),
         "Richieste e Documenti", [
             "  Richiesta certificati e relazioni terapeutiche",
             "  Selezione tipologia documento",
             "  Motivazione e note aggiuntive",
             "  Tracciamento stato richiesta in tempo reale",
             "  Download documento quando pronto",
             "  Storico richieste effettuate",
         ], SECONDARY)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.2),
         "Reclami e Comunicazioni", [
             "  Invio reclami con titolo e descrizione",
             "  Segnalazione assenza del paziente",
             "  Ricezione notifiche dal centro",
             "  Accesso multi-paziente per famiglie",
             "  Interfaccia semplice e intuitiva",
             "  Supporto iOS e Android",
         ], PURPLE)

add_footer(slide, 13, TOTAL_SLIDES)


# =============================================================================
# SLIDE 14 - PANNELLO SUPER ADMIN
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), PURPLE)
add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5),
            "PANNELLO SUPER ADMIN", font_size=14, color=RGBColor(0xDD, 0xD6, 0xFE))
add_textbox(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
            "Controllo Totale del Sistema", font_size=26, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.4),
         "Gestione Ruoli e Permessi", [
             "  Creazione e modifica ruoli personalizzati",
             "  Assegnazione permessi granulari per ruolo",
             "  Visualizzazione matrice ruoli-permessi",
             "  Export configurazione permessi",
         ], PURPLE)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(2.4),
         "Gestione Utenti", [
             "  Panoramica tutti gli utenti del sistema",
             "  Assegnazione permessi individuali",
             "  Gestione accessi per piattaforma",
             "  Attivazione/disattivazione utenti",
         ], SECONDARY)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(2.4),
         "Distretti Sanitari (ASL)", [
             "  Configurazione distretti sul territorio",
             "  Riferimento ASL per ogni distretto",
             "  Associazione pazienti-distretto",
             "  Gestione multi-sede e multi-ASL",
         ], GREEN)

add_card(slide, Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.4),
         "Log Attivita (Audit Trail)", [
             "  Registro completo di ogni azione",
             "  Chi ha fatto cosa e quando",
             "  Tracciamento modifiche ai dati",
             "  Filtri per utente, data e tipo azione",
         ], ORANGE)

add_card(slide, Inches(4.6), Inches(4.2), Inches(3.8), Inches(2.4),
         "Impostazioni di Sistema", [
             "  Configurazione parametri globali",
             "  Impostazioni calendario (range date)",
             "  Parametri sicurezza e autenticazione",
             "  Configurazione notifiche",
         ], DARK)

add_card(slide, Inches(8.7), Inches(4.2), Inches(4.1), Inches(2.4),
         "Documentazione API", [
             "  Swagger UI integrata per test API",
             "  Documentazione interattiva endpoints",
             "  Test autenticazione JWT dall'interfaccia",
             "  Esempi request/response in tempo reale",
         ], RED)

add_footer(slide, 14, TOTAL_SLIDES)


# =============================================================================
# SLIDE 15 - SICUREZZA E SPECIFICHE TECNICHE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
            "Sicurezza e Architettura", font_size=30, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.4),
         "Autenticazione e Accesso", [
             "  JWT (JSON Web Token) con firma crittografica",
             "  Access Token (1 ora) + Refresh Token (30 giorni)",
             "  Scadenza password ogni 90 giorni",
             "  Login biometrico (impronta, Face ID)",
         ], RED)

add_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(2.4),
         "Protezione Dati", [
             "  RBAC (Role-Based Access Control) granulare",
             "  Protezione CSRF su tutte le form",
             "  Validazione input server-side completa",
             "  Audit trail di ogni operazione",
         ], ORANGE)

add_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(2.4),
         "Infrastruttura", [
             "  Architettura web moderna e scalabile",
             "  Database relazionale MySQL ottimizzato",
             "  API RESTful con documentazione Swagger",
             "  Caching intelligente per performance",
         ], SECONDARY)

add_card(slide, Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.4),
         "App Mobile", [
             "  React Native: iOS e Android nativi",
             "  Comunicazione sicura HTTPS",
             "  Storage locale crittografato",
             "  Notifiche push in tempo reale",
         ], GREEN)

add_card(slide, Inches(4.6), Inches(4.2), Inches(3.8), Inches(2.4),
         "Conformita e Standard", [
             "  Gestione dati sensibili conforme GDPR",
             "  Password policy robusta con caratteri speciali",
             "  Sessioni con timeout automatico",
             "  Separazione ambienti (dev/staging/prod)",
         ], PURPLE)

add_card(slide, Inches(8.7), Inches(4.2), Inches(4.1), Inches(2.4),
         "Scalabilita", [
             "  Supporto multi-sede nativo",
             "  Gestione multi-distretto ASL",
             "  Performance ottimizzate con viste materializzate",
             "  Architettura modulare ed estensibile",
         ], DARK)

add_footer(slide, 15, TOTAL_SLIDES)


# =============================================================================
# SLIDE 16 - CHIUSURA / CTA
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, PRIMARY)

# Decorazioni
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), ACCENT)
add_shape_bg(slide, Inches(0), SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), ACCENT)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, SLIDE_W - Inches(4), Inches(-1), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x25, 0x4B, 0xBF)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), SLIDE_H - Inches(3), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x25, 0x4B, 0xBF)
circle2.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "TherapyCRM", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5.5), Inches(2.8), Inches(2.5), Inches(0.05), ACCENT)

add_textbox(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.8),
            "La piattaforma che trasforma la gestione del vostro centro terapeutico",
            font_size=22, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)

# Punti chiave finali
benefits = [
    "Gestione centralizzata di pazienti, terapisti e appuntamenti",
    "App mobile per terapisti e famiglie sempre connessi",
    "Statistiche e report per decisioni basate sui dati",
    "Sicurezza e conformita ai massimi standard",
    "Scalabile: da una sede a molte, da 10 a 1000 pazienti",
]

txBox = slide.shapes.add_textbox(Inches(2.5), Inches(4.0), Inches(8), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, b in enumerate(benefits):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"  {b}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
    p.font.name = 'Calibri'
    p.space_before = Pt(6)
    p.space_after = Pt(4)
    p.alignment = PP_ALIGN.CENTER

add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
            "Contattateci per una demo personalizzata",
            font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# Salva
output_path = "/Applications/XAMPP/xamppfiles/htdocs/TherapyCRM/TherapyCRM_Presentazione.pptx"
prs.save(output_path)
print(f"Presentazione salvata in: {output_path}")
print(f"Totale slide: {len(prs.slides)}")
